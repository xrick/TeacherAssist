# backend/app/services/pptx_analyzer.py
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import logging

logger = logging.getLogger(__name__)


class ImageInfo:
    """Image metadata extracted from slide"""
    def __init__(
        self,
        image_bytes: bytes,
        left: int,
        top: int,
        width: int,
        height: int,
        image_format: str = "unknown"
    ):
        self.image_bytes = image_bytes
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.image_format = image_format

    def to_dict(self) -> Dict[str, Any]:
        return {
            "left": self.left,
            "top": self.top,
            "width": self.width,
            "height": self.height,
            "format": self.image_format,
            "size_bytes": len(self.image_bytes)
        }


class SlideData:
    """Structured slide data extracted from PPTX"""
    def __init__(
        self,
        index: int,
        title: str = "",
        content: List[str] = None,
        layout: str = "",
        has_image: bool = False,
        images: List[ImageInfo] = None,
        speaker_notes: str = ""
    ):
        self.index = index
        self.title = title
        self.content = content or []
        self.layout = layout
        self.has_image = has_image
        self.images = images or []
        self.speaker_notes = speaker_notes

    def to_dict(self) -> Dict[str, Any]:
        return {
            "index": self.index,
            "title": self.title,
            "content": self.content,
            "layout": self.layout,
            "has_image": self.has_image,
            "images": [img.to_dict() for img in self.images],
            "speaker_notes": self.speaker_notes
        }


class PythonPptxAnalyzer:
    """PPTX parser for extracting structured data from Presenton-generated presentations"""

    def __init__(self):
        self.logger = logger

    def parse_presentation(self, pptx_bytes: bytes) -> List[SlideData]:
        """
        Parse PPTX bytes and extract structured slide data

        Args:
            pptx_bytes: PPTX file content as bytes

        Returns:
            List of SlideData objects containing extracted information

        Raises:
            ValueError: If PPTX is invalid or cannot be parsed
        """
        try:
            pptx_stream = io.BytesIO(pptx_bytes)
            presentation = Presentation(pptx_stream)

            slides_data = []
            for idx, slide in enumerate(presentation.slides):
                slide_data = self._parse_slide(idx, slide)
                slides_data.append(slide_data)

            self.logger.info(f"Successfully parsed {len(slides_data)} slides from PPTX")
            return slides_data

        except Exception as e:
            self.logger.error(f"Failed to parse PPTX: {str(e)}")
            raise ValueError(f"Invalid PPTX file: {str(e)}")

    def _parse_slide(self, index: int, slide) -> SlideData:
        """Parse individual slide and extract all components"""
        title = ""
        content = []
        images = []
        speaker_notes = ""

        # Extract layout name
        layout_name = slide.slide_layout.name if hasattr(slide, 'slide_layout') else "Unknown"

        # Extract title (try multiple methods)
        # Method 1: slide.shapes.title
        if hasattr(slide.shapes, 'title') and slide.shapes.title is not None:
            title = self._extract_text_from_shape(slide.shapes.title)

        # Method 2: Find placeholder type 1 (title)
        if not title:
            for shape in slide.shapes:
                if shape.has_text_frame and hasattr(shape, 'is_placeholder'):
                    if shape.is_placeholder and shape.placeholder_format.type == 1:
                        title = self._extract_text_from_shape(shape)
                        break

        # Extract content (non-title text)
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = self._extract_text_from_shape(shape)
                # Skip title and empty text
                if text and text != title:
                    content.append(text)

        # Extract images
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_info = self._extract_image(shape)
                    if image_info:
                        images.append(image_info)
                except Exception as e:
                    self.logger.warning(f"Failed to extract image from slide {index}: {e}")

        # Extract speaker notes
        if hasattr(slide, 'notes_slide') and slide.notes_slide:
            try:
                notes_text_frame = slide.notes_slide.notes_text_frame
                if notes_text_frame:
                    speaker_notes = notes_text_frame.text.strip()
            except Exception as e:
                self.logger.warning(f"Failed to extract speaker notes from slide {index}: {e}")

        return SlideData(
            index=index,
            title=title,
            content=content,
            layout=layout_name,
            has_image=len(images) > 0,
            images=images,
            speaker_notes=speaker_notes
        )

    def _extract_text_from_shape(self, shape) -> str:
        """
        Extract text from shape's text frame

        Args:
            shape: python-pptx Shape object

        Returns:
            Extracted text as string
        """
        if not shape.has_text_frame:
            return ""

        text_parts = []
        for paragraph in shape.text_frame.paragraphs:
            paragraph_text = paragraph.text.strip()
            if paragraph_text:
                text_parts.append(paragraph_text)

        return "\n".join(text_parts)

    def _extract_image(self, shape) -> Optional[ImageInfo]:
        """
        Extract image binary data and metadata from picture shape

        Args:
            shape: python-pptx Picture shape

        Returns:
            ImageInfo object with image data and position metadata
        """
        try:
            # Get image binary data
            image = shape.image
            image_bytes = image.blob

            # Get position and size (in EMUs - English Metric Units)
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            # Detect image format
            image_format = image.ext if hasattr(image, 'ext') else "unknown"

            return ImageInfo(
                image_bytes=image_bytes,
                left=left,
                top=top,
                width=width,
                height=height,
                image_format=image_format
            )

        except Exception as e:
            self.logger.error(f"Failed to extract image: {e}")
            return None

    def detect_layout_type(self, slide_data: SlideData) -> str:
        """
        Detect semantic layout type from slide structure

        Args:
            slide_data: Parsed slide data

        Returns:
            Layout type category (title, content, image, two_column, etc.)
        """
        # Title slide: has title, minimal content, no images
        if slide_data.title and len(slide_data.content) <= 1 and not slide_data.has_image:
            return "title"

        # Image slide: has images
        if slide_data.has_image:
            if len(slide_data.images) > 1:
                return "image_gallery"
            return "image"

        # Two-column heuristic: many content items (>5)
        if len(slide_data.content) > 5:
            return "two_column"

        # Standard content slide
        if slide_data.title and slide_data.content:
            return "content"

        # Blank or unknown
        return "blank"

    def get_summary(self, slides_data: List[SlideData]) -> Dict[str, Any]:
        """
        Generate summary statistics from parsed slides

        Args:
            slides_data: List of parsed slide data

        Returns:
            Summary dictionary with statistics
        """
        total_slides = len(slides_data)
        total_images = sum(len(slide.images) for slide in slides_data)
        slides_with_images = sum(1 for slide in slides_data if slide.has_image)
        slides_with_notes = sum(1 for slide in slides_data if slide.speaker_notes)

        layout_counts = {}
        for slide in slides_data:
            layout = slide.layout
            layout_counts[layout] = layout_counts.get(layout, 0) + 1

        return {
            "total_slides": total_slides,
            "total_images": total_images,
            "slides_with_images": slides_with_images,
            "slides_with_notes": slides_with_notes,
            "layout_distribution": layout_counts,
            "avg_content_per_slide": sum(len(s.content) for s in slides_data) / total_slides if total_slides > 0 else 0
        }
