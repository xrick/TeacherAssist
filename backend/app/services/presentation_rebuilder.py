# backend/app/services/presentation_rebuilder.py
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io

from app.services.content_improver import ImprovedSlideData
from app.services.pptx_analyzer import ImageInfo

logger = logging.getLogger(__name__)


class PresentationRebuilder:
    """Rebuild presentation using templates with improved content"""

    DEFAULT_FONT = "Microsoft JhengHei"
    TITLE_SIZE = Pt(40)
    CONTENT_SIZE = Pt(20)
    LINE_SPACING = 1.2
    TITLE_COLOR = RGBColor(31, 73, 125)
    CONTENT_COLOR = RGBColor(0, 0, 0)

    def __init__(self, template_dir: Optional[Path] = None):
        self.template_dir = template_dir or Path(__file__).parent.parent.parent.parent / "refData" / "free_templates"
        self.logger = logger

    def rebuild_presentation(
        self,
        slides: List[ImprovedSlideData],
        template_path: Optional[str] = None,
        output_format: str = "bytes"
    ) -> bytes:
        """
        Rebuild presentation from improved slides using template

        Args:
            slides: List of improved slide data
            template_path: Path to template PPTX (or filename if in template_dir)
            output_format: Output format ("bytes" or "file")

        Returns:
            PPTX file as bytes

        Raises:
            ValueError: If template loading fails or rebuild fails
        """
        try:
            prs = self._load_template(template_path)

            self._clear_existing_slides(prs)

            for slide_data in slides:
                self._add_slide_from_data(prs, slide_data)

            self._apply_global_formatting(prs)

            pptx_stream = io.BytesIO()
            prs.save(pptx_stream)
            pptx_bytes = pptx_stream.getvalue()

            self.logger.info(f"Successfully rebuilt presentation with {len(slides)} slides")
            return pptx_bytes

        except Exception as e:
            self.logger.error(f"Failed to rebuild presentation: {e}")
            raise ValueError(f"Presentation rebuild failed: {str(e)}")

    def _load_template(self, template_path: Optional[str]) -> Presentation:
        """
        Load template presentation

        Args:
            template_path: Full path or filename in template_dir

        Returns:
            Presentation object
        """
        if not template_path:
            self.logger.info("No template specified, using blank presentation")
            return Presentation()

        if Path(template_path).exists():
            full_path = Path(template_path)
        else:
            full_path = self.template_dir / template_path

        if not full_path.exists():
            self.logger.warning(f"Template not found: {full_path}, using blank presentation")
            return Presentation()

        try:
            prs = Presentation(str(full_path))
            self.logger.info(f"Loaded template: {full_path.name} ({len(prs.slide_layouts)} layouts)")
            return prs
        except Exception as e:
            self.logger.error(f"Failed to load template {full_path}: {e}")
            raise ValueError(f"Invalid template file: {str(e)}")

    def _clear_existing_slides(self, prs: Presentation):
        """Remove all existing slides from presentation"""
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

    def _add_slide_from_data(self, prs: Presentation, slide_data: ImprovedSlideData):
        """
        Add slide to presentation from slide data

        Args:
            prs: Presentation object
            slide_data: Slide data to add
        """
        layout = self._smart_select_layout(prs, slide_data)

        slide = prs.slides.add_slide(layout)

        title_text = slide_data.improved_title or slide_data.title
        content = slide_data.improved_content or slide_data.content

        if slide.shapes.title:
            self._set_title(slide.shapes.title, title_text)

        if content and len(content) > 0:
            self._add_content(slide, content, layout)

        if slide_data.has_image and slide_data.images:
            try:
                self._restore_images(slide, slide_data.images)
            except Exception as e:
                self.logger.warning(f"Failed to restore images for slide {slide_data.index}: {e}")

    def _smart_select_layout(
        self,
        prs: Presentation,
        slide_data: ImprovedSlideData
    ) -> Any:
        """
        Intelligently select layout based on slide data

        Args:
            prs: Presentation object
            slide_data: Slide data

        Returns:
            Selected slide layout
        """
        layouts = prs.slide_layouts

        if slide_data.index == 0:
            return self._find_layout_by_type(layouts, "title")

        if slide_data.has_image:
            image_layout = self._find_layout_by_type(layouts, "picture")
            if image_layout:
                return image_layout

        content = slide_data.improved_content or slide_data.content
        if content and len(content) > 5:
            two_col_layout = self._find_layout_by_type(layouts, "two_column")
            if two_col_layout:
                return two_col_layout

        content_layout = self._find_layout_by_type(layouts, "content")
        if content_layout:
            return content_layout

        return layouts[1] if len(layouts) > 1 else layouts[0]

    def _find_layout_by_type(self, layouts, layout_type: str) -> Optional[Any]:
        """
        Find layout by type keyword

        Args:
            layouts: Available slide layouts
            layout_type: Type keyword (title, content, picture, two_column, blank)

        Returns:
            Matching layout or None
        """
        type_keywords = {
            "title": ["title", "標題"],
            "content": ["content", "內容", "body"],
            "picture": ["picture", "image", "圖片", "photo"],
            "two_column": ["two", "column", "雙欄", "comparison"],
            "blank": ["blank", "空白"]
        }

        keywords = type_keywords.get(layout_type, [layout_type])

        for layout in layouts:
            layout_name_lower = layout.name.lower()
            if any(keyword.lower() in layout_name_lower for keyword in keywords):
                return layout

        if layout_type == "title" and len(layouts) > 0:
            return layouts[0]

        if layout_type == "content" and len(layouts) > 1:
            return layouts[1]

        return None

    def _set_title(self, title_shape, text: str):
        """
        Set title with formatting

        Args:
            title_shape: Title shape object
            text: Title text
        """
        title_shape.text = text

        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = self.DEFAULT_FONT
                run.font.size = self.TITLE_SIZE
                run.font.bold = True
                run.font.color.rgb = self.TITLE_COLOR

    def _add_content(self, slide, content_list: List[str], layout):
        """
        Add content to slide

        Args:
            slide: Slide object
            content_list: List of content points
            layout: Slide layout
        """
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:
                content_placeholder = shape
                break

        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()

            for i, content in enumerate(content_list):
                if "[過渡句]" in content:
                    continue

                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = content
                p.level = 0
                p.space_before = Pt(6)
                p.space_after = Pt(6)

                for run in p.runs:
                    run.font.name = self.DEFAULT_FONT
                    run.font.size = self.CONTENT_SIZE
                    run.font.color.rgb = self.CONTENT_COLOR
        else:
            self._create_text_box(slide, content_list)

    def _create_text_box(self, slide, content_list: List[str]):
        """
        Create text box for content when placeholder not available

        Args:
            slide: Slide object
            content_list: List of content points
        """
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for i, content in enumerate(content_list):
            if "[過渡句]" in content:
                continue

            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = f"• {content}"
            p.space_before = Pt(6)
            p.space_after = Pt(6)

            for run in p.runs:
                run.font.name = self.DEFAULT_FONT
                run.font.size = Pt(18)

    def _restore_images(self, slide, images: List[ImageInfo]):
        """
        Restore images to slide with exact positions

        Args:
            slide: Slide object
            images: List of ImageInfo objects

        NOTE: This is a high-risk operation (as identified in estimation report)
        Uses careful binary handling and exact position coordinates
        """
        for image_info in images:
            try:
                image_stream = io.BytesIO(image_info.image_bytes)

                left = image_info.left
                top = image_info.top
                width = image_info.width
                height = image_info.height

                slide.shapes.add_picture(
                    image_stream,
                    left=left,
                    top=top,
                    width=width,
                    height=height
                )

                self.logger.debug(f"Restored image at ({left}, {top}) size ({width}, {height})")

            except Exception as e:
                self.logger.error(f"Failed to restore individual image: {e}")
                continue

    def _apply_global_formatting(self, prs: Presentation):
        """
        Apply global formatting polish to all slides

        Args:
            prs: Presentation object
        """
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self._polish_text_frame(shape.text_frame)

    def _polish_text_frame(self, text_frame):
        """
        Polish text frame formatting

        Args:
            text_frame: TextFrame object
        """
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        for paragraph in text_frame.paragraphs:
            paragraph.line_spacing = self.LINE_SPACING

            for run in paragraph.runs:
                if not run.font.name or run.font.name == 'Calibri':
                    run.font.name = self.DEFAULT_FONT

                if not run.font.size or run.font.size < Pt(14):
                    run.font.size = Pt(18)

    def get_template_info(self, template_path: str) -> Dict[str, Any]:
        """
        Get template metadata

        Args:
            template_path: Path to template file

        Returns:
            Template information dict
        """
        try:
            prs = self._load_template(template_path)

            layouts = []
            for idx, layout in enumerate(prs.slide_layouts):
                layout_info = {
                    "index": idx,
                    "name": layout.name,
                    "placeholders": []
                }

                for shape in layout.placeholders:
                    placeholder_info = {
                        "idx": shape.placeholder_format.idx,
                        "type": shape.placeholder_format.type,
                        "name": shape.name
                    }
                    layout_info["placeholders"].append(placeholder_info)

                layouts.append(layout_info)

            return {
                "template_name": Path(template_path).name,
                "layout_count": len(prs.slide_layouts),
                "layouts": layouts
            }

        except Exception as e:
            self.logger.error(f"Failed to get template info: {e}")
            return {
                "template_name": template_path,
                "error": str(e)
            }
