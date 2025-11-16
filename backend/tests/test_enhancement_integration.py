# backend/tests/test_enhancement_integration.py
import pytest
import sys
import os
from io import BytesIO
from pptx import Presentation

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from app.services.pptx_analyzer import PythonPptxAnalyzer
from app.services.content_improver import ContentImprover
from app.services.presentation_rebuilder import PresentationRebuilder
from app.services.template_manager import TemplateManager


class TestEnhancementIntegration:
    """Integration tests for enhancement pipeline"""

    @pytest.fixture
    def sample_pptx_bytes(self):
        """Create sample PPTX for testing"""
        prs = Presentation()

        title_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_layout)
        slide1.shapes.title.text = "測試簡報"

        content_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(content_layout)
        slide2.shapes.title.text = "內容頁"
        body = slide2.placeholders[1]
        tf = body.text_frame
        tf.text = "要點一"
        tf.add_paragraph().text = "要點二"

        pptx_stream = BytesIO()
        prs.save(pptx_stream)
        return pptx_stream.getvalue()

    @pytest.mark.asyncio
    async def test_full_enhancement_pipeline_without_template(self, sample_pptx_bytes):
        """Test complete enhancement pipeline without template"""
        analyzer = PythonPptxAnalyzer()
        improver = ContentImprover()
        rebuilder = PresentationRebuilder()

        slides_data = analyzer.parse_presentation(sample_pptx_bytes)
        assert len(slides_data) == 2

        improved_slides = await improver.improve_slides_batch(slides_data, "測試簡報")
        assert len(improved_slides) == 2

        enhanced_pptx_bytes = rebuilder.rebuild_presentation(improved_slides, template_path=None)
        assert isinstance(enhanced_pptx_bytes, bytes)
        assert len(enhanced_pptx_bytes) > 0

        prs = Presentation(BytesIO(enhanced_pptx_bytes))
        assert len(prs.slides) == 2

    @pytest.mark.asyncio
    async def test_full_enhancement_pipeline_with_template(self, sample_pptx_bytes):
        """Test complete enhancement pipeline with template"""
        analyzer = PythonPptxAnalyzer()
        improver = ContentImprover()
        rebuilder = PresentationRebuilder()
        template_manager = TemplateManager()

        slides_data = analyzer.parse_presentation(sample_pptx_bytes)
        improved_slides = await improver.improve_slides_batch(slides_data, "測試簡報")

        template_name = template_manager.get_recommended_template("educational")

        if template_name:
            template_path = template_manager.templates_dir / template_name
            enhanced_pptx_bytes = rebuilder.rebuild_presentation(
                improved_slides,
                template_path=str(template_path)
            )

            assert isinstance(enhanced_pptx_bytes, bytes)
            prs = Presentation(BytesIO(enhanced_pptx_bytes))
            assert len(prs.slides) == 2

    def test_template_manager_recommendations(self):
        """Test template recommendations"""
        template_manager = TemplateManager()

        # Templates may not exist in test environment, so test graceful fallback
        templates = template_manager.list_templates()

        if templates:
            # If templates exist, test recommendations
            educational = template_manager.get_recommended_template("educational")
            assert educational is not None

            thesis = template_manager.get_recommended_template("thesis")
            assert thesis is not None

            general = template_manager.get_recommended_template("general")
            assert general is not None
        else:
            # If no templates, ensure graceful handling
            educational = template_manager.get_recommended_template("educational")
            # Should not crash, may return None
            assert educational is None or isinstance(educational, str)

    @pytest.mark.asyncio
    async def test_improvement_quality(self, sample_pptx_bytes):
        """Test that improvement increases content quality"""
        analyzer = PythonPptxAnalyzer()
        improver = ContentImprover()

        slides_data = analyzer.parse_presentation(sample_pptx_bytes)
        improved_slides = await improver.improve_slides_batch(slides_data, "測試簡報")

        for original, improved in zip(slides_data, improved_slides):
            assert improved.improved_title is not None

            if improved.improvement_applied:
                assert len(improved.improved_content) >= 3


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
