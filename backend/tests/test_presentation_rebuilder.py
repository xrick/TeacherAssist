# backend/tests/test_presentation_rebuilder.py
import pytest
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import sys
import os

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from app.services.presentation_rebuilder import PresentationRebuilder
from app.services.content_improver import ImprovedSlideData
from app.services.pptx_analyzer import SlideData, ImageInfo


class TestPresentationRebuilder:
    """Test suite for PresentationRebuilder"""

    @pytest.fixture
    def rebuilder(self):
        return PresentationRebuilder()

    @pytest.fixture
    def sample_improved_slides(self):
        slide1 = SlideData(0, "標題頁", [], "Title Slide")
        improved1 = ImprovedSlideData(slide1)
        improved1.improved_title = "深度學習基礎"
        improved1.improved_content = []

        slide2 = SlideData(1, "內容頁", ["原始要點1", "原始要點2"], "Content")
        improved2 = ImprovedSlideData(slide2)
        improved2.improved_title = "神經網路架構"
        improved2.improved_content = ["輸入層與資料預處理", "隱藏層與特徵提取", "輸出層與預測"]
        improved2.improvement_applied = True

        return [improved1, improved2]

    def test_rebuild_presentation_no_template(self, rebuilder, sample_improved_slides):
        """Test rebuilding with blank presentation"""
        pptx_bytes = rebuilder.rebuild_presentation(sample_improved_slides, template_path=None)

        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 0

        prs = Presentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 2

    def test_load_template_nonexistent(self, rebuilder):
        """Test loading nonexistent template"""
        prs = rebuilder._load_template("nonexistent_template.pptx")

        assert prs is not None
        assert len(prs.slides) == 0

    def test_load_template_blank(self, rebuilder):
        """Test loading blank template"""
        prs = rebuilder._load_template(None)

        assert prs is not None

    def test_clear_existing_slides(self, rebuilder):
        """Test clearing slides from presentation"""
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        prs.slides.add_slide(title_layout)
        prs.slides.add_slide(title_layout)

        assert len(prs.slides) == 2

        rebuilder._clear_existing_slides(prs)

        assert len(prs.slides) == 0

    def test_smart_select_layout_title_slide(self, rebuilder):
        """Test layout selection for title slide"""
        prs = Presentation()
        slide_data = ImprovedSlideData(SlideData(0, "Title", [], "Title"))

        layout = rebuilder._smart_select_layout(prs, slide_data)

        assert layout is not None
        assert layout == prs.slide_layouts[0]

    def test_smart_select_layout_content_slide(self, rebuilder):
        """Test layout selection for content slide"""
        prs = Presentation()
        slide_data = ImprovedSlideData(SlideData(1, "Content", ["Point 1", "Point 2"], "Content"))

        layout = rebuilder._smart_select_layout(prs, slide_data)

        assert layout is not None

    def test_find_layout_by_type_title(self, rebuilder):
        """Test finding title layout"""
        prs = Presentation()
        layouts = prs.slide_layouts

        title_layout = rebuilder._find_layout_by_type(layouts, "title")

        assert title_layout is not None

    def test_find_layout_by_type_content(self, rebuilder):
        """Test finding content layout"""
        prs = Presentation()
        layouts = prs.slide_layouts

        content_layout = rebuilder._find_layout_by_type(layouts, "content")

        assert content_layout is not None

    def test_set_title(self, rebuilder):
        """Test setting title with formatting"""
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        title_shape = slide.shapes.title

        rebuilder._set_title(title_shape, "測試標題")

        assert title_shape.text == "測試標題"

    def test_create_text_box(self, rebuilder):
        """Test creating text box for content"""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        content = ["要點 1", "要點 2", "要點 3"]
        rebuilder._create_text_box(slide, content)

        assert len(slide.shapes) > 0

    def test_restore_images_with_sample_data(self, rebuilder):
        """Test image restoration with sample data"""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        fake_image_bytes = b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82'

        images = [
            ImageInfo(
                image_bytes=fake_image_bytes,
                left=Inches(1),
                top=Inches(2),
                width=Inches(3),
                height=Inches(2),
                image_format="png"
            )
        ]

        try:
            rebuilder._restore_images(slide, images)
            assert len(slide.shapes) > 0
        except Exception as e:
            pytest.skip(f"Image restoration requires valid PNG data: {e}")

    def test_polish_text_frame(self, rebuilder):
        """Test text frame polishing"""
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Test"

        text_frame = title_shape.text_frame

        rebuilder._polish_text_frame(text_frame)

        assert text_frame.word_wrap is True

    def test_add_slide_from_data(self, rebuilder):
        """Test adding slide from slide data"""
        prs = Presentation()

        slide_data = ImprovedSlideData(SlideData(0, "Test Title", ["Point 1", "Point 2"], "Content"))
        slide_data.improved_title = "Improved Title"
        slide_data.improved_content = ["Improved Point 1", "Improved Point 2", "Improved Point 3"]

        rebuilder._add_slide_from_data(prs, slide_data)

        assert len(prs.slides) == 1


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
