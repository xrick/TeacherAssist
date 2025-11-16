# backend/tests/test_pptx_analyzer.py
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import sys
import os

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from app.services.pptx_analyzer import PythonPptxAnalyzer, SlideData, ImageInfo


class TestPythonPptxAnalyzer:
    """Test suite for PythonPptxAnalyzer"""

    @pytest.fixture
    def analyzer(self):
        return PythonPptxAnalyzer()

    @pytest.fixture
    def simple_pptx_bytes(self):
        """Create a simple PPTX for testing"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "測試標題"
        subtitle.text = "測試副標題"

        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = "內容頁標題"
        tf = body.text_frame
        tf.text = "第一個要點"
        p = tf.add_paragraph()
        p.text = "第二個要點"
        p.level = 0

        pptx_stream = BytesIO()
        prs.save(pptx_stream)
        return pptx_stream.getvalue()

    def test_parse_presentation_success(self, analyzer, simple_pptx_bytes):
        """Test successful PPTX parsing"""
        slides_data = analyzer.parse_presentation(simple_pptx_bytes)

        assert len(slides_data) == 2
        assert all(isinstance(slide, SlideData) for slide in slides_data)

    def test_parse_presentation_invalid_bytes(self, analyzer):
        """Test parsing with invalid PPTX bytes"""
        invalid_bytes = b"not a valid pptx file"

        with pytest.raises(ValueError, match="Invalid PPTX file"):
            analyzer.parse_presentation(invalid_bytes)

    def test_extract_title(self, analyzer, simple_pptx_bytes):
        """Test title extraction"""
        slides_data = analyzer.parse_presentation(simple_pptx_bytes)

        assert slides_data[0].title == "測試標題"
        assert slides_data[1].title == "內容頁標題"

    def test_extract_content(self, analyzer, simple_pptx_bytes):
        """Test content extraction"""
        slides_data = analyzer.parse_presentation(simple_pptx_bytes)

        slide_1_content = slides_data[1].content
        assert len(slide_1_content) > 0
        assert any("第一個要點" in content or "第二個要點" in content for content in slide_1_content)

    def test_detect_layout_type(self, analyzer):
        """Test layout type detection"""
        title_slide = SlideData(
            index=0,
            title="Title Only",
            content=[],
            has_image=False
        )
        assert analyzer.detect_layout_type(title_slide) == "title"

        content_slide = SlideData(
            index=1,
            title="Content Title",
            content=["Point 1", "Point 2", "Point 3"],
            has_image=False
        )
        assert analyzer.detect_layout_type(content_slide) == "content"

        image_slide = SlideData(
            index=2,
            title="Image Title",
            content=[],
            has_image=True,
            images=[ImageInfo(b"fake", 0, 0, 100, 100)]
        )
        assert analyzer.detect_layout_type(image_slide) == "image"

        two_column_slide = SlideData(
            index=3,
            title="Two Column",
            content=["1", "2", "3", "4", "5", "6"],
            has_image=False
        )
        assert analyzer.detect_layout_type(two_column_slide) == "two_column"

    def test_get_summary(self, analyzer, simple_pptx_bytes):
        """Test summary generation"""
        slides_data = analyzer.parse_presentation(simple_pptx_bytes)
        summary = analyzer.get_summary(slides_data)

        assert summary["total_slides"] == 2
        assert "layout_distribution" in summary
        assert "avg_content_per_slide" in summary
        assert summary["total_images"] == 0
        assert summary["slides_with_images"] == 0

    def test_slide_data_to_dict(self):
        """Test SlideData serialization"""
        slide = SlideData(
            index=0,
            title="Test Title",
            content=["Point 1", "Point 2"],
            layout="Title Slide",
            has_image=False,
            speaker_notes="Test notes"
        )

        data_dict = slide.to_dict()

        assert data_dict["index"] == 0
        assert data_dict["title"] == "Test Title"
        assert len(data_dict["content"]) == 2
        assert data_dict["layout"] == "Title Slide"
        assert data_dict["has_image"] is False
        assert data_dict["speaker_notes"] == "Test notes"

    def test_image_info_to_dict(self):
        """Test ImageInfo serialization"""
        image = ImageInfo(
            image_bytes=b"fake_image_data",
            left=100,
            top=200,
            width=300,
            height=400,
            image_format="png"
        )

        data_dict = image.to_dict()

        assert data_dict["left"] == 100
        assert data_dict["top"] == 200
        assert data_dict["width"] == 300
        assert data_dict["height"] == 400
        assert data_dict["format"] == "png"
        assert data_dict["size_bytes"] == len(b"fake_image_data")


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
