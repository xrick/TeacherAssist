#!/usr/bin/env python3
"""
Template Validation Script
Validates refData/free_templates PPTX files for python-pptx compatibility
"""
import sys
import os
from pathlib import Path
from pptx import Presentation
from typing import Dict, List, Tuple

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '../..')))

TEMPLATES_DIR = Path(__file__).parent.parent.parent / "refData" / "free_templates"


class TemplateValidator:
    """Validates PPTX templates for compatibility"""

    def __init__(self, templates_dir: Path):
        self.templates_dir = templates_dir
        self.results = []

    def validate_all(self) -> List[Dict]:
        """Validate all templates in directory"""
        if not self.templates_dir.exists():
            print(f"❌ Templates directory not found: {self.templates_dir}")
            return []

        template_files = list(self.templates_dir.glob("*.pptx"))
        print(f"Found {len(template_files)} template files\n")

        for template_file in template_files:
            result = self.validate_template(template_file)
            self.results.append(result)

        return self.results

    def validate_template(self, template_path: Path) -> Dict:
        """Validate single template file"""
        result = {
            "filename": template_path.name,
            "path": str(template_path),
            "valid": False,
            "error": None,
            "slide_count": 0,
            "layout_count": 0,
            "layouts": [],
            "has_title_layout": False,
            "has_content_layout": False
        }

        print(f"Validating: {template_path.name}")

        try:
            prs = Presentation(str(template_path))

            result["valid"] = True
            result["slide_count"] = len(prs.slides)
            result["layout_count"] = len(prs.slide_layouts)

            for idx, layout in enumerate(prs.slide_layouts):
                layout_info = {
                    "index": idx,
                    "name": layout.name,
                    "placeholders": []
                }

                for shape in layout.placeholders:
                    placeholder_info = {
                        "idx": shape.placeholder_format.idx,
                        "type": shape.placeholder_format.type,
                        "name": shape.name
                    }
                    layout_info["placeholders"].append(placeholder_info)

                result["layouts"].append(layout_info)

                layout_name_lower = layout.name.lower()
                if "title" in layout_name_lower and idx < 2:
                    result["has_title_layout"] = True
                if "content" in layout_name_lower or "body" in layout_name_lower:
                    result["has_content_layout"] = True

            status = "✅ Valid"
            print(f"  {status}")
            print(f"  - Slides: {result['slide_count']}")
            print(f"  - Layouts: {result['layout_count']}")
            print(f"  - Has title layout: {result['has_title_layout']}")
            print(f"  - Has content layout: {result['has_content_layout']}")

        except Exception as e:
            result["error"] = str(e)
            print(f"  ❌ Invalid: {str(e)}")

        print()
        return result

    def generate_report(self) -> str:
        """Generate validation report"""
        total = len(self.results)
        valid = sum(1 for r in self.results if r["valid"])
        invalid = total - valid

        report_lines = [
            "=" * 60,
            "Template Validation Report",
            "=" * 60,
            f"Total templates: {total}",
            f"Valid: {valid} ({valid/total*100:.1f}%)" if total > 0 else "Valid: 0",
            f"Invalid: {invalid}",
            "",
            "=" * 60,
            "Valid Templates:",
            "=" * 60
        ]

        for result in self.results:
            if result["valid"]:
                report_lines.append(f"\n✅ {result['filename']}")
                report_lines.append(f"   Slides: {result['slide_count']}")
                report_lines.append(f"   Layouts: {result['layout_count']}")
                report_lines.append(f"   Available layouts:")
                for layout in result["layouts"]:
                    placeholders_count = len(layout["placeholders"])
                    report_lines.append(f"     [{layout['index']}] {layout['name']} ({placeholders_count} placeholders)")

        if invalid > 0:
            report_lines.extend([
                "",
                "=" * 60,
                "Invalid Templates:",
                "=" * 60
            ])
            for result in self.results:
                if not result["valid"]:
                    report_lines.append(f"\n❌ {result['filename']}")
                    report_lines.append(f"   Error: {result['error']}")

        report_lines.extend([
            "",
            "=" * 60,
            "Recommendation:",
            "=" * 60
        ])

        if valid >= 6:
            report_lines.append("✅ Sufficient templates available (≥6). Proceed with implementation.")
        elif valid >= 4:
            report_lines.append("⚠️  Limited templates (4-5). Consider MVP with verified templates only.")
        else:
            report_lines.append("❌ Insufficient valid templates (<4). Review template compatibility.")

        return "\n".join(report_lines)

    def save_report(self, output_path: Path):
        """Save report to file"""
        report = self.generate_report()
        output_path.write_text(report, encoding="utf-8")
        print(f"\n📄 Report saved to: {output_path}")


def main():
    print("Template Validation Script")
    print("=" * 60)
    print(f"Templates directory: {TEMPLATES_DIR}\n")

    validator = TemplateValidator(TEMPLATES_DIR)
    results = validator.validate_all()

    print(validator.generate_report())

    output_path = Path(__file__).parent.parent / "claudedocs" / "template_validation_report.md"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    validator.save_report(output_path)

    compatible_count = sum(1 for r in results if r["valid"])
    total_count = len(results)

    print(f"\n{'='*60}")
    print(f"Compatibility: {compatible_count}/{total_count} ({compatible_count/total_count*100:.1f}%)" if total_count > 0 else "No templates found")
    print(f"{'='*60}")

    sys.exit(0 if compatible_count >= 6 else 1)


if __name__ == "__main__":
    main()
